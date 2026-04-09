import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

path = 'c:/Users/greco/.cursor/projects/apex-drone-solutions/AG Drones NJ_Pitch Deck.pptx'
prs = Presentation(path)

# ============================================================
# SLIDE 1: Title - Update headline numbers
# ============================================================
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace('52.6%', '62.1%')
                run.text = run.text.replace('62x ROI', '67x ROI')

# ============================================================
# SLIDE 3: Solution - T25 -> T50, update specs
# ============================================================
slide3 = prs.slides[2]
for shape in slide3.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace('DJI T25', 'DJI Agras T50')
                run.text = run.text.replace('20 acres/hour', '45+ acres/hour')

# ============================================================
# SLIDE 5: Revenue Projections - Update CAGR
# ============================================================
slide5 = prs.slides[4]
for shape in slide5.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace('47.6% CAGR', '62.7% CAGR')

# ============================================================
# SLIDE 6: Profitability Path - Update all numbers
# ============================================================
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            # Cards
            run.text = run.text.replace('$1.43M', '$1.69M')
            run.text = run.text.replace('52.6%', '62.1%')
            run.text = run.text.replace('$998K', '$1.18M')
            # Gross margins
            run.text = run.text.replace('54.4%', '70.5%')
            run.text = run.text.replace('59.3%', '71.8%')
            run.text = run.text.replace('61.3%', '72.7%')
            run.text = run.text.replace('64.0%', '73.1%')
            run.text = run.text.replace('63.9%', '72.7%')
            # EBITDA row
            run.text = run.text.replace('$96K', '$165K')
            run.text = run.text.replace('$293K', '$398K')
            run.text = run.text.replace('$652K', '$826K')
            run.text = run.text.replace('$997K', '$1.19M')
            # EBITDA margins
            run.text = run.text.replace('24.7%', '42.4%')
            run.text = run.text.replace('38.1%', '51.8%')
            run.text = run.text.replace('46.4%', '58.8%')
            run.text = run.text.replace('50.5%', '60.5%')
            # Net income
            run.text = run.text.replace('$53K', '$98K')
            run.text = run.text.replace('$196K', '$268K')
            run.text = run.text.replace('$451K', '$572K')
            run.text = run.text.replace('$688K', '$825K')
            # Net margins
            run.text = run.text.replace('13.6%', '25.2%')
            run.text = run.text.replace('25.5%', '34.9%')
            run.text = run.text.replace('32.1%', '40.7%')
            run.text = run.text.replace('34.8%', '41.8%')
            run.text = run.text.replace('36.6%', '43.3%')
            # Cumulative
            run.text = run.text.replace('$2.39M', '$2.94M')
            run.text = run.text.replace('$1.95M', '$2.42M')

# ============================================================
# SLIDE 7: Competitive Advantages - Add per-acre model
# ============================================================
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if 'WOTC Veteran Hiring Tax Credits' in run.text:
                    run.text = 'Variable Per-Acre Cost Structure + WOTC Credits'
                if '$4,800' in run.text and 'annual savings' in run.text:
                    run.text = 'Pilots and crew paid per-acre ($3.70/acre), making largest cost 100% variable with revenue. Plus $4,800-$9,600/yr WOTC veteran hiring tax credits.'

# ============================================================
# SLIDE 8: Unit Economics - Update all numbers
# ============================================================
slide8 = prs.slides[7]
for shape in slide8.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = run.text.replace('$16.00', '$16.21')
            run.text = run.text.replace('Y5 $14.50', 'Y5 $15.14')
            run.text = run.text.replace('$7.38', '$4.78')
            run.text = run.text.replace('Y5 $5.46', 'Y5 $4.13')
            run.text = run.text.replace('$8.82', '$11.43')
            run.text = run.text.replace('Y5 $9.68', 'Y5 $11.01')
            run.text = run.text.replace('Y1 $194K', 'Y1 $195K')
            run.text = run.text.replace('$239K', '$282K')
            run.text = run.text.replace('Y1 $48K', 'Y1 $82K')
            run.text = run.text.replace('Y5 $239K', 'Y5 $282K')

# ============================================================
# SLIDE 10: Investment & Exit - Update all
# ============================================================
slide10 = prs.slides[9]
for shape in slide10.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = run.text.replace('62x Total ROI', '67x Total ROI')
            run.text = run.text.replace('$164K initial', '$179K initial')
            run.text = run.text.replace('$164,200', '$178,700')
            run.text = run.text.replace('$85,000', '$100,000')
            run.text = run.text.replace('$3.0M', '$3.56M')
            run.text = run.text.replace('18.3x', '19.9x')
            run.text = run.text.replace('$5.74M', '$6.77M')
            run.text = run.text.replace('$7.17M', '$8.46M')
            run.text = run.text.replace('$8.60M', '$10.16M')
            run.text = run.text.replace('$10.18M', '$12.03M')

# ============================================================
# SLIDE 11: The Ask - Update amount
# ============================================================
slide11 = prs.slides[10]
for shape in slide11.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            run.text = run.text.replace('$164,200', '$178,700')
            run.text = run.text.replace('62x Total Return', '67x Total Return')

# ============================================================
# ADD NEW SLIDE: Regulatory Moat & Timeline
# ============================================================
layout = prs.slide_layouts[1]  # Title and Content
new_slide = prs.slides.add_slide(layout)

title_shape = new_slide.placeholders[0]
title_shape.text = 'Regulatory Moat & Timeline'

content_shape = new_slide.placeholders[1]
tf = content_shape.text_frame
tf.clear()

lines = [
    ('7-step FAA certification creates a durable barrier to entry', True),
    ('', False),
    ('1. FAA Part 107 Commercial Drone License - Written exam, $175, pass rate 99%+', False),
    ('2. Form LLC & Register Business - State filing + EIN for drone operations', False),
    ('3. Register Drones with FAA - Physical paperwork required for 55lb+ aircraft', False),
    ('4. 3rd Class Medical Exam - FAA-authorized physician (AME), similar to CDL', False),
    ('5. 44807 Exemption - FAA petition for 55lb+ drones (longest step: 2-4 months)', False),
    ('6. NJ Chemical Applicator License - CORE + AERIAL via Rutgers 1-day program', False),
    ('7. FAA Part 137 Certificate - Agricultural aircraft operator certification', False),
    ('', False),
    ('Timeline: 3-6 months from start to first paid acre', False),
    ('Cost: ~$10K in licensing, legal, and medical fees', False),
    ('Result: Multi-layered regulatory moat that casual operators cannot replicate', False),
]

for i, (text, bold) in enumerate(lines):
    if i == 0:
        para = tf.paragraphs[0]
    else:
        para = tf.add_paragraph()
    para.text = text
    para.space_after = Pt(4)
    if para.runs:
        run = para.runs[0]
        run.font.size = Pt(14)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

# Move new slide from last position to after slide 7 (index 7)
slide_list = prs.slides._sldIdLst
slide_ids = list(slide_list)
new_slide_elem = slide_ids[-1]
slide_list.remove(new_slide_elem)
slide_list.insert(7, new_slide_elem)

prs.save(path)
print("Pitch deck updated successfully!")
print("  - Slide 1: EBITDA margin 52.6% -> 62.1%, ROI 62x -> 67x")
print("  - Slide 3: DJI T25 -> T50, 20 acres/hr -> 45+ acres/hr")
print("  - Slide 5: CAGR 47.6% -> 62.7%")
print("  - Slide 6: All profitability numbers updated")
print("  - Slide 7: Added per-acre cost structure to competitive advantages")
print("  - NEW Slide 8: Regulatory Moat & Timeline (7-step FAA process)")
print("  - Slide 9: Unit economics updated (COGS/acre $7.38->$4.78)")
print("  - Slide 11: Investment/exit values updated")
print("  - Slide 12: The Ask updated to $178,700")
